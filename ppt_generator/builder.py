"""
PPT Builder - constructs PowerPoint slides from analyzed content.
"""

from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .theme import COLORS, SLIDE_W, SLIDE_H, MARGIN, CONTENT_W


class PPTBuilder:
    """Build PowerPoint slides from analyzed content."""

    # Font defaults — CJK for Chinese/Japanese/Korean, Calibri for Latin
    CJK_FONT = "Microsoft YaHei"
    LATIN_FONT = "Calibri"

    def __init__(self, font: str = ""):
        """
        Args:
            font: Default font name. If empty, auto-detects based on first content.
        """
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._explicit_font = font  # user-specified font takes priority

    def _detect_font(self, text: str = "") -> str:
        """Choose font based on content language."""
        if self._explicit_font:
            return self._explicit_font
        if text:
            cjk = sum(1 for c in text if '\u4e00' <= c <= '\u9fff' or '\u3040' <= c <= '\u30ff')
            latin = sum(1 for c in text if c.isascii() and c.isalpha())
            return self.CJK_FONT if cjk > latin else self.LATIN_FONT
        return self.LATIN_FONT

    def _blank_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rect(self, slide, left, top, width, height, color, rounded=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _text(self, slide, left, top, width, height, text,
              size=18, color=None, bold=False, align=PP_ALIGN.LEFT,
              font=""):
        font = font or self._detect_font(text)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["text"]
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        return box

    def _add_accent_bar(self, slide, y=Inches(0), color=None):
        self._rect(slide, Inches(0), y, SLIDE_W, Inches(0.06),
                   color or COLORS["primary"])

    def _add_slide_number(self, slide, num, total):
        self._text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.3),
                   f"{num} / {total}", size=10, color=COLORS["slate400"],
                   align=PP_ALIGN.RIGHT)

    # ─── Slide types ────────────────────────────────────────────

    def slide_title(self, chapter_title: str, book_title: str,
                    chapter_number: int, page_count: int):
        slide = self._blank_slide()
        # Dark background
        self._rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["dark"])
        # Accent bar
        self._rect(slide, 0, 0, SLIDE_W, Inches(0.08), COLORS["primary_light"])

        # Chapter number badge
        if chapter_number:
            self._rect(slide, MARGIN, Inches(2.0),
                       Inches(1.8), Inches(0.5), COLORS["primary"], rounded=True)
            self._text(slide, MARGIN, Inches(2.05), Inches(1.8), Inches(0.4),
                       f"Chapter {chapter_number}", size=16,
                       color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)

        # Title
        top = Inches(2.8) if chapter_number else Inches(2.2)
        self._text(slide, MARGIN, top, Inches(11), Inches(1.8),
                   chapter_title, size=44, color=COLORS["white"], bold=True)

        # Subtitle
        if book_title:
            self._text(slide, MARGIN, Inches(5.0), Inches(8), Inches(0.5),
                       book_title, size=20, color=COLORS["slate400"])

        # Page count
        self._text(slide, MARGIN, Inches(6.5), Inches(4), Inches(0.3),
                   f"{page_count} pages", size=14, color=COLORS["slate400"])

    def slide_outline(self, sections: list):
        slide = self._blank_slide()
        self._add_accent_bar(slide)

        self._text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.8),
                   "Outline", size=32, color=COLORS["dark"], bold=True)

        self._rect(slide, MARGIN, Inches(1.1), Inches(1.5), Inches(0.04),
                   COLORS["primary"])

        # Section list
        y = Inches(1.5)
        for i, sec in enumerate(sections):
            title = getattr(sec, 'title', str(sec))
            if not title:
                continue
            # Number
            self._rect(slide, MARGIN, y, Inches(0.45), Inches(0.45),
                       COLORS["primary"], rounded=True)
            self._text(slide, MARGIN, y + Inches(0.05), Inches(0.45), Inches(0.35),
                       str(i + 1), size=16, color=COLORS["white"],
                       bold=True, align=PP_ALIGN.CENTER)

            # Title
            self._text(slide, Inches(1.5), y + Inches(0.05), Inches(10), Inches(0.35),
                       title, size=20, color=COLORS["text"])

            y += Inches(0.6)
            if y > Inches(6.5):
                break

    def slide_section_divider(self, section_title: str, section_num: int,
                              total_sections: int):
        from pptx.dml.color import RGBColor
        slide = self._blank_slide()
        self._rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["primary"])

        # Section number
        self._text(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(0.6),
                   f"Section {section_num}", size=18,
                   color=RGBColor(0xBF, 0xDB, 0xFE))

        # Title
        self._text(slide, MARGIN, Inches(2.8), CONTENT_W, Inches(2.0),
                   section_title, size=40, color=COLORS["white"], bold=True)

        # Progress bar
        progress = section_num / total_sections
        bar_w = Inches(4)
        self._rect(slide, MARGIN, Inches(5.5), bar_w, Inches(0.15),
                   RGBColor(0x1E, 0x40, 0xAF), rounded=True)
        self._rect(slide, MARGIN, Inches(5.5),
                   Inches(4 * progress), Inches(0.15),
                   COLORS["white"], rounded=True)

    def slide_content(self, title: str, points: list[str],
                      source_pages: str = "", subtitle: str = ""):
        slide = self._blank_slide()
        self._add_accent_bar(slide)

        # Title
        self._text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.8),
                   title, size=28, color=COLORS["dark"], bold=True)

        # Subtitle if provided
        title_bottom = Inches(1.1)
        if subtitle:
            self._text(slide, MARGIN, Inches(1.1), CONTENT_W, Inches(0.5),
                       subtitle, size=16, color=COLORS["text_secondary"])
            title_bottom = Inches(1.5)

        self._rect(slide, MARGIN, title_bottom, Inches(1.5), Inches(0.04),
                   COLORS["primary"])

        # Bullet points as cards
        y = title_bottom + Inches(0.4)
        for i, point in enumerate(points):
            # Card background
            card_h = Inches(0.8) if len(point) < 100 else Inches(1.1)
            self._rect(slide, MARGIN, y, CONTENT_W, card_h,
                       COLORS["bg_card"], rounded=True)

            # Bullet marker
            self._rect(slide, Inches(1.0), y + Inches(0.15),
                       Inches(0.08), Inches(0.08), COLORS["primary"], rounded=True)

            # Text
            self._text(slide, Inches(1.3), y + Inches(0.1),
                       Inches(10.5), card_h - Inches(0.2),
                       point, size=16, color=COLORS["text_secondary"])

            y += card_h + Inches(0.15)
            if y > Inches(6.2):
                break

        # Source reference
        if source_pages:
            self._text(slide, Inches(10), Inches(7.0), Inches(3), Inches(0.3),
                       source_pages, size=10, color=COLORS["slate400"],
                       align=PP_ALIGN.RIGHT)

    def slide_content_with_image(self, title: str, points: list[str],
                                  image_path: Union[str, Path],
                                  source_pages: str = "",
                                  image_caption: str = ""):
        """Content slide with key points on the left and a diagram/image on the right."""
        image_path = Path(image_path)
        slide = self._blank_slide()
        self._add_accent_bar(slide)

        # Title (full width)
        self._text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.8),
                   title, size=28, color=COLORS["dark"], bold=True)
        self._rect(slide, MARGIN, Inches(1.1), Inches(1.5), Inches(0.04),
                   COLORS["primary"])

        # Left side: key points (narrower)
        text_w = Inches(6.0)
        y = Inches(1.5)
        for point in points[:4]:
            card_h = Inches(0.7) if len(point) < 100 else Inches(0.95)
            self._rect(slide, MARGIN, y, text_w, card_h,
                       COLORS["bg_card"], rounded=True)
            self._rect(slide, Inches(1.0), y + Inches(0.15),
                       Inches(0.08), Inches(0.08), COLORS["primary"], rounded=True)
            self._text(slide, Inches(1.3), y + Inches(0.08),
                       Inches(5.5), card_h - Inches(0.15),
                       point, size=14, color=COLORS["text_secondary"])
            y += card_h + Inches(0.1)
            if y > Inches(6.0):
                break

        # Right side: image
        try:
            img_x = Inches(7.2)
            img_max_w = Inches(5.5)
            img_max_h = Inches(5.0)
            pic = slide.shapes.add_picture(str(image_path),
                                           img_x, Inches(1.5),
                                           img_max_w, img_max_h)
            aspect = pic.image.size[0] / pic.image.size[1]
            if aspect > (5.5 / 5.0):
                pic.width = img_max_w
                pic.height = int(img_max_w / aspect)
            else:
                pic.height = img_max_h
                pic.width = int(img_max_h * aspect)
            # Center vertically within image area
            pic.top = int(Inches(1.5) + (img_max_h - pic.height) / 2)
        except Exception:
            self._text(slide, Inches(7.5), Inches(3.5), Inches(5), Inches(1),
                       f"[Image: {image_path.name}]", size=14,
                       color=COLORS["slate400"])

        # Image caption
        if image_caption:
            self._text(slide, Inches(7.2), Inches(6.8), Inches(5.5), Inches(0.3),
                       image_caption, size=10, color=COLORS["slate400"],
                       align=PP_ALIGN.CENTER)

        # Source reference
        if source_pages:
            self._text(slide, MARGIN, Inches(7.0), Inches(3), Inches(0.3),
                       source_pages, size=10, color=COLORS["slate400"],
                       align=PP_ALIGN.LEFT)

    def slide_quote(self, quote: str, attribution: str = "",
                    context: str = ""):
        slide = self._blank_slide()

        # Left accent bar
        self._rect(slide, MARGIN, Inches(1.5), Inches(0.08), Inches(4.5),
                   COLORS["primary"])

        # Quote text
        self._text(slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(3.5),
                   f'"{quote}"', size=24, color=COLORS["text"])

        # Attribution
        if attribution:
            self._text(slide, Inches(1.2), Inches(5.5), Inches(8), Inches(0.5),
                       f"— {attribution}", size=16,
                       color=COLORS["text_secondary"])

        # Context
        if context:
            self._text(slide, Inches(1.2), Inches(6.2), Inches(8), Inches(0.4),
                       context, size=13, color=COLORS["slate400"])

    def slide_definition(self, term: str, definition: str, context: str = ""):
        slide = self._blank_slide()
        self._add_accent_bar(slide, color=COLORS["accent"])

        self._text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
                   "Key Concept", size=14, color=COLORS["accent"], bold=True)

        # Term
        self._text(slide, MARGIN, Inches(1.2), CONTENT_W, Inches(1.0),
                   term, size=36, color=COLORS["dark"], bold=True)

        # Definition card
        self._rect(slide, MARGIN, Inches(2.5), CONTENT_W, Inches(2.5),
                   COLORS["accent_light"], rounded=True)
        self._text(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(2.0),
                   definition, size=20, color=COLORS["text"])

        if context:
            self._text(slide, MARGIN, Inches(5.5), CONTENT_W, Inches(1.0),
                       context, size=16, color=COLORS["text_secondary"])

    def slide_two_column(self, title: str, left_title: str, left_points: list[str],
                         right_title: str, right_points: list[str]):
        # Deduplicate points
        left_points = list(dict.fromkeys(left_points))[:4]
        right_points = list(dict.fromkeys(right_points))[:4]
        slide = self._blank_slide()
        self._add_accent_bar(slide)

        self._text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.8),
                   title, size=28, color=COLORS["dark"], bold=True)
        self._rect(slide, MARGIN, Inches(1.1), Inches(1.5), Inches(0.04),
                   COLORS["primary"])

        col_w = Inches(5.5)
        left_x = MARGIN
        right_x = Inches(7)

        # Left column
        self._rect(slide, left_x, Inches(1.5), col_w, Inches(0.5),
                   COLORS["primary"], rounded=True)
        self._text(slide, left_x + Inches(0.2), Inches(1.55), col_w, Inches(0.4),
                   left_title, size=16, color=COLORS["white"], bold=True)

        y = Inches(2.2)
        for point in left_points[:4]:
            self._text(slide, left_x + Inches(0.2), y, Inches(5), Inches(0.7),
                       f"• {point}", size=14, color=COLORS["text_secondary"])
            y += Inches(0.8)

        # Right column
        self._rect(slide, right_x, Inches(1.5), col_w, Inches(0.5),
                   COLORS["accent"], rounded=True)
        self._text(slide, right_x + Inches(0.2), Inches(1.55), col_w, Inches(0.4),
                   right_title, size=16, color=COLORS["white"], bold=True)

        y = Inches(2.2)
        for point in right_points[:4]:
            self._text(slide, right_x + Inches(0.2), y, Inches(5), Inches(0.7),
                       f"• {point}", size=14, color=COLORS["text_secondary"])
            y += Inches(0.8)

    def slide_image(self, image_path: Union[str, Path], caption: str,
                    source: str = ""):
        image_path = Path(image_path)
        slide = self._blank_slide()
        self._add_accent_bar(slide)

        self._text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
                   caption, size=24, color=COLORS["dark"], bold=True)

        try:
            max_w = Inches(10)
            max_h = Inches(5)
            pic = slide.shapes.add_picture(str(image_path),
                                           Inches(1.5), Inches(1.2),
                                           max_w, max_h)
            aspect = pic.image.size[0] / pic.image.size[1]
            if aspect > (10 / 5):
                pic.width = max_w
                pic.height = int(max_w / aspect)
            else:
                pic.height = max_h
                pic.width = int(max_h * aspect)
            pic.left = int((SLIDE_W - pic.width) / 2)
        except Exception:
            self._text(slide, Inches(3), Inches(3), Inches(7), Inches(1),
                       f"[Image: {image_path.name}]", size=16,
                       color=COLORS["slate400"])

        if source:
            self._text(slide, MARGIN, Inches(6.8), CONTENT_W, Inches(0.3),
                       source, size=12, color=COLORS["slate400"],
                       align=PP_ALIGN.CENTER)

    def slide_summary(self, chapter_title: str, section_titles: list[str],
                      stats: dict):
        from pptx.dml.color import RGBColor
        slide = self._blank_slide()
        self._rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["dark"])
        self._rect(slide, 0, 0, SLIDE_W, Inches(0.08), COLORS["primary_light"])

        self._text(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
                   "Summary", size=36, color=COLORS["white"], bold=True)
        self._text(slide, MARGIN, Inches(1.3), CONTENT_W, Inches(0.5),
                   chapter_title, size=18, color=COLORS["slate400"])

        # Stats row
        stats_items = [
            ("Pages", str(stats.get("pages", 0))),
            ("Sections", str(stats.get("sections", 0))),
            ("Slides", str(stats.get("slides", 0))),
            ("Images", str(stats.get("images", 0))),
        ]

        card_w = Inches(2.5)
        gap = Inches(0.4)
        start_x = MARGIN
        y = Inches(2.3)

        for i, (label, value) in enumerate(stats_items):
            cx = start_x + i * (card_w + gap)
            self._rect(slide, cx, y, card_w, Inches(1.5),
                       COLORS["slate800"], rounded=True)
            self._text(slide, cx, y + Inches(0.2), card_w, Inches(0.8),
                       value, size=40, color=COLORS["primary_light"],
                       bold=True, align=PP_ALIGN.CENTER)
            self._text(slide, cx, y + Inches(1.0), card_w, Inches(0.3),
                       label, size=14, color=COLORS["slate400"],
                       align=PP_ALIGN.CENTER)

        # Key sections
        y2 = Inches(4.3)
        self._text(slide, MARGIN, y2, Inches(3), Inches(0.4),
                   "Key Sections", size=16, color=COLORS["slate400"], bold=True)
        y2 += Inches(0.5)

        for title in section_titles[:6]:
            self._text(slide, Inches(1.2), y2, Inches(10), Inches(0.35),
                       f"▸ {title}", size=15, color=RGBColor(0xCB, 0xD5, 0xE1))
            y2 += Inches(0.4)
